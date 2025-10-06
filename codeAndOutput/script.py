

# ======= Script starts here =======
import os
import io
import sys
import math
import warnings
warnings.filterwarnings('ignore')

# Install dependencies not always present
try:
    import pptx
except Exception:
    print('Installing python-pptx...')
    !pip install python-pptx

try:
    import statsmodels.api as sm
except Exception:
    print('Installing statsmodels...')
    !pip install statsmodels

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# For Colab interactive uploads
try:
    from google.colab import files
    COLAB = True
except Exception:
    COLAB = False

WORKDIR = Path('/content')
os.makedirs(WORKDIR, exist_ok=True)

# Expected filenames (common variants)
expected_files = {
    'flights': ['Flight Level Data.csv','Flight_Level_Data.csv','Flight Level Data - Sheet1.csv'],
    'pnr': ['PNR+Flight+Level+Data.csv','PNR_Flight_Level_Data.csv','PNR Flight Level Data.csv'],
    'pnr_remark': ['PNR Remark Level Data.csv','PNR_Remark_Level_Data.csv','PNR Remark Level Data.csv'],
    'bags': ['Bag+Level+Data.csv','Bag Level Data.csv','Bag_Level_Data.csv'],
    'airports': ['Airports Data.csv','Airports_Data.csv']
}

found = {}
for key, names in expected_files.items():
    for n in names:
        p = WORKDIR / n
        if p.exists():
            found[key] = p
            break

# If not all files found and running in Colab, prompt upload
missing = [k for k in expected_files.keys() if k not in found]
if missing and COLAB:
    print('Some expected CSVs are missing. Please upload the following files when prompted:')
    print('\n'.join(missing))
    uploaded = files.upload()
    # save uploaded
    for fn, content in uploaded.items():
        target = WORKDIR / fn
        with open(target, 'wb') as f:
            f.write(content)
        print('Saved', target)
    # try to find again
    for key, names in expected_files.items():
        for n in names:
            p = WORKDIR / n
            if p.exists():
                found[key] = p
                break

# If still missing, attempt to search directory
for key, names in expected_files.items():
    if key not in found:
        for p in WORKDIR.iterdir():
            low = p.name.lower()
            for n in names:
                if n.lower() in low:
                    found[key] = p
                    break
            if key in found:
                break

# Ensure we have all
missing = [k for k in expected_files.keys() if k not in found]
if missing:
    print('ERROR: Missing files:', missing)
    print('Please upload the required CSVs to /content and re-run.')
    raise SystemExit(1)

print('Using files:')
for k,v in found.items():
    print(k, '->', v)

# Read CSVs
parse_dates_flights = ['scheduled_departure_datetime_local','scheduled_arrival_datetime_local','actual_departure_datetime_local','actual_arrival_datetime_local']
flights = pd.read_csv(found['flights'], parse_dates=[c for c in parse_dates_flights if c in pd.read_csv(found['flights'], nrows=0).columns], low_memory=False)
pnr = pd.read_csv(found['pnr'], parse_dates=['pnr_creation_date'] if 'pnr_creation_date' in pd.read_csv(found['pnr'], nrows=0).columns else [], low_memory=False)
pnr_remark = pd.read_csv(found['pnr_remark'], parse_dates=['pnr_creation_date'] if 'pnr_creation_date' in pd.read_csv(found['pnr_remark'], nrows=0).columns else [], low_memory=False)
bags = pd.read_csv(found['bags'], parse_dates=['bag_tag_issue_date'] if 'bag_tag_issue_date' in pd.read_csv(found['bags'], nrows=0).columns else [], low_memory=False)
airports = pd.read_csv(found['airports'], low_memory=False)

# Normalize column names
for df in [flights, pnr, pnr_remark, bags, airports]:
    df.columns = [c.strip() for c in df.columns]

# === EDA & Feature engineering ===
# 1. departure_delay_minutes
if all(c in flights.columns for c in ['actual_departure_datetime_local','scheduled_departure_datetime_local']):
    flights['departure_delay_minutes'] = (pd.to_datetime(flights['actual_departure_datetime_local']) - pd.to_datetime(flights['scheduled_departure_datetime_local'])).dt.total_seconds()/60.0
else:
    flights['departure_delay_minutes'] = np.nan

# 2. ground time deficit
if 'scheduled_ground_time_minutes' in flights.columns and 'minimum_turn_minutes' in flights.columns:
    flights['ground_time_deficit'] = flights['scheduled_ground_time_minutes'] - flights['minimum_turn_minutes']
else:
    flights['ground_time_deficit'] = np.nan

# 3. bag summaries
bags['bag_type'] = bags['bag_type'].astype(str)
bags['bag_type_clean'] = bags['bag_type'].str.lower()
bag_summary = bags.pivot_table(index=['company_id','flight_number','scheduled_departure_date_local','scheduled_departure_station_code'],
                               columns='bag_type_clean', values='bag_tag_unique_number', aggfunc='count', fill_value=0).reset_index()
# find likely checked/transfer columns
checked_col = next((c for c in bag_summary.columns if 'check' in str(c).lower()), None)
transfer_col = next((c for c in bag_summary.columns if 'trans' in str(c).lower()), None)
if checked_col:
    bag_summary = bag_summary.rename(columns={checked_col:'checked_bags'})
else:
    bag_summary['checked_bags']=0
if transfer_col:
    bag_summary = bag_summary.rename(columns={transfer_col:'transfer_bags'})
else:
    bag_summary['transfer_bags']=0
bag_summary['transfer_to_checked_ratio'] = bag_summary.apply(lambda r: r['transfer_bags']/r['checked_bags'] if r['checked_bags']>0 else np.nan, axis=1)

# 4. PNR aggregation
# detect basic economy column name variations
pnr_cols = [c.lower() for c in pnr.columns]
if 'basic_economy_ind' in pnr_cols:
    basic_col = [c for c in pnr.columns if c.lower()=='basic_economy_ind'][0]
elif 'basic_economy_pax' in pnr_cols:
    basic_col = [c for c in pnr.columns if c.lower()=='basic_economy_pax'][0]
else:
    basic_col = None

agg_map = {'total_pax':'sum', 'lap_child_count':'sum'}
# ensure columns exist
agg_cols = {}
if 'total_pax' in pnr.columns:
    agg_cols['total_pax'] = ('total_pax','sum')
if 'lap_child_count' in pnr.columns:
    agg_cols['lap_child_count'] = ('lap_child_count','sum')
if basic_col:
    agg_cols['basic_economy_pax'] = (basic_col,'sum')

if not agg_cols:
    print('PNR file does not contain expected columns for aggregation. Continuing with defaults (zeros).')
    pnr_flight = pd.DataFrame(columns=['company_id','flight_number','scheduled_departure_date_local','scheduled_departure_station_code','total_pax_flight','lap_child_count','basic_economy_pax'])
else:
    pnr_flight = pnr.groupby(['company_id','flight_number','scheduled_departure_date_local','scheduled_departure_station_code']).agg(**{k:v for k,v in agg_cols.items()}).reset_index()
    # rename to consistent names
    rename_map = {}
    if 'total_pax' in agg_cols:
        rename_map['total_pax'] = 'total_pax_flight'
    if 'lap_child_count' in agg_cols:
        rename_map['lap_child_count'] = 'lap_child_count'
    if 'basic_economy_pax' in agg_cols:
        rename_map['basic_economy_pax'] = 'basic_economy_pax'
    pnr_flight = pnr_flight.rename(columns=rename_map)

# 5. SSR aggregation: join pnr_remark -> pnr to attach scheduled date
if 'record_locator' in pnr_remark.columns and 'record_locator' in pnr.columns:
    pnr_key = pnr[['record_locator','flight_number','scheduled_departure_date_local']].drop_duplicates()
    pnr_remark_joined = pnr_remark.merge(pnr_key, on=['record_locator','flight_number'], how='left')
    pnr_remark_joined['special_service_request'] = pnr_remark_joined['special_service_request'].fillna('')
    ssr_count = pnr_remark_joined.groupby(['flight_number','scheduled_departure_date_local'])['special_service_request'].count().reset_index().rename(columns={'special_service_request':'ssr_count'})
else:
    # fallback: if pnr_remark has scheduled_departure_date_local join directly
    if 'scheduled_departure_date_local' in pnr_remark.columns:
        pnr_remark['special_service_request'] = pnr_remark['special_service_request'].fillna('')
        ssr_count = pnr_remark.groupby(['flight_number','scheduled_departure_date_local'])['special_service_request'].count().reset_index().rename(columns={'special_service_request':'ssr_count'})
    else:
        ssr_count = pd.DataFrame(columns=['flight_number','scheduled_departure_date_local','ssr_count'])

# Merge to master
flights['scheduled_departure_date_local'] = flights['scheduled_departure_date_local'].astype(str)
if not pnr_flight.empty:
    pnr_flight['scheduled_departure_date_local'] = pnr_flight['scheduled_departure_date_local'].astype(str)
bag_summary['scheduled_departure_date_local'] = bag_summary['scheduled_departure_date_local'].astype(str)

master = flights.merge(pnr_flight, on=['company_id','flight_number','scheduled_departure_date_local','scheduled_departure_station_code'], how='left')\
                .merge(bag_summary[['company_id','flight_number','scheduled_departure_date_local','scheduled_departure_station_code','checked_bags','transfer_bags','transfer_to_checked_ratio']], 
                       on=['company_id','flight_number','scheduled_departure_date_local','scheduled_departure_station_code'], how='left')

master = master.merge(ssr_count, on=['flight_number','scheduled_departure_date_local'], how='left')
master['ssr_count'] = master['ssr_count'].fillna(0)

# Fill NaNs with sensible defaults
for col in ['total_pax_flight','checked_bags','transfer_bags','transfer_to_checked_ratio','basic_economy_pax','lap_child_count']:
    if col in master.columns:
        master[col] = master[col].fillna(0)

master['scheduled_ground_time_minutes'] = master.get('scheduled_ground_time_minutes', master.get('minimum_turn_minutes'))
master['scheduled_ground_time_minutes'] = master['scheduled_ground_time_minutes'].fillna(master.get('minimum_turn_minutes', np.nan))
master['ground_time_deficit'] = master['scheduled_ground_time_minutes'] - master.get('minimum_turn_minutes', master['scheduled_ground_time_minutes'])
master['passenger_load_pct'] = master.apply(lambda r: r['total_pax_flight']/r['total_seats'] if ('total_pax_flight' in r.index and r['total_seats']>0) else 0, axis=1)

# Hot transfer detection
bags['is_hot_transfer'] = bags['bag_type'].str.contains('hot', case=False, na=False)
hot_counts = bags.groupby(['company_id','flight_number','scheduled_departure_date_local'])['is_hot_transfer'].sum().reset_index().rename(columns={'is_hot_transfer':'hot_transfer_bags'})
master = master.merge(hot_counts, on=['company_id','flight_number','scheduled_departure_date_local'], how='left')
master['hot_transfer_bags'] = master['hot_transfer_bags'].fillna(0)

# departure delay fill
master['departure_delay_minutes'] = master['departure_delay_minutes'].fillna(0)

# --- Correlations ---
feature_cols = ['ground_time_deficit','total_pax_flight','passenger_load_pct','checked_bags','transfer_bags','transfer_to_checked_ratio','ssr_count','hot_transfer_bags','basic_economy_pax','lap_child_count']
corrs = {}
for c in feature_cols:
    if c in master.columns:
        try:
            cor = master[c].corr(master['departure_delay_minutes'])
            corrs[c]=cor if not np.isnan(cor) else 0.0
        except Exception:
            corrs[c]=0.0
    else:
        corrs[c]=0.0

# --- Build difficulty score ---
# robust scale function
def robust_scale(series):
    med = series.median()
    iqr = series.quantile(0.75) - series.quantile(0.25)
    if iqr == 0 or np.isnan(iqr):
        iqr = series.std() if series.std()>0 else 1.0
    return (series - med) / iqr

score_df = master.copy()
# weights from abs correlation with floor
weights = {}
total_w = 0.0
for f in feature_cols:
    if f in score_df.columns:
        w = abs(corrs.get(f,0.0))
        if w == 0:
            w = 0.05
        weights[f]=w
        total_w += w
# normalize
for k in weights:
    weights[k] = weights[k]/total_w if total_w>0 else 1.0/len(weights)

# scaled features
for f in weights.keys():
    score_df[f+'_scaled'] = robust_scale(score_df[f].fillna(0))
    if f == 'ground_time_deficit':
        score_df[f+'_scaled'] = -score_df[f+'_scaled']  # less ground time -> more difficult

# difficulty raw
score_df['difficulty_score_raw'] = 0.0
for f,w in weights.items():
    score_df['difficulty_score_raw'] += score_df[f+'_scaled'] * w

# per-day normalization to 0-100
score_df['scheduled_departure_date_local_day'] = score_df['scheduled_departure_date_local']

def normalize_group(g):
    vals = g['difficulty_score_raw']
    if vals.std()==0 or np.isnan(vals.std()):
        return pd.Series(np.zeros(len(vals)), index=g.index)
    z = (vals - vals.mean()) / vals.std()
    mm = (z - z.min()) / (z.max() - z.min()) if (z.max()-z.min())!=0 else z*0
    return mm*100

score_df['difficulty_score'] = score_df.groupby('scheduled_departure_date_local_day', group_keys=False).apply(normalize_group).reset_index(drop=True)
score_df['daily_rank'] = score_df.groupby('scheduled_departure_date_local_day')['difficulty_score'].rank(method='first', ascending=False)
score_df['daily_pct_rank'] = score_df.groupby('scheduled_departure_date_local_day')['difficulty_score'].rank(pct=True, ascending=False)

# classification
def classify_pct(pct):
    if pct <= 0.2:
        return 'Difficult'
    elif pct <= 0.8:
        return 'Medium'
    else:
        return 'Easy'
score_df['difficulty_class'] = score_df['daily_pct_rank'].apply(classify_pct)

# Save CSV
out_csv = WORKDIR / 'test_yourname.csv'
cols_out = [c for c in ['company_id','flight_number','scheduled_departure_date_local','scheduled_departure_station_code','scheduled_departure_datetime_local','scheduled_ground_time_minutes','minimum_turn_minutes','ground_time_deficit','total_seats','total_pax_flight','passenger_load_pct','checked_bags','transfer_bags','transfer_to_checked_ratio','hot_transfer_bags','ssr_count','departure_delay_minutes','difficulty_score','daily_rank','daily_pct_rank','difficulty_class'] if c in score_df.columns]
score_df[cols_out].to_csv(out_csv, index=False)
print('Saved output CSV to', out_csv)

# === Visuals (matplotlib) ===
plot_dir = WORKDIR / 'plots'
plot_dir.mkdir(exist_ok=True)

# 1. Histogram of difficulty_score
plt.figure(figsize=(8,4))
plt.hist(score_df['difficulty_score'].dropna(), bins=40)
plt.title('Difficulty Score Distribution (0-100)')
plt.xlabel('Difficulty score')
plt.ylabel('Count')
plt.tight_layout()
plt.savefig(plot_dir / 'difficulty_hist.png')
plt.close()

# 2. Scatter: difficulty vs departure delay
plt.figure(figsize=(8,4))
plt.scatter(score_df['difficulty_score'], score_df['departure_delay_minutes'], alpha=0.4)
plt.xlabel('Difficulty score')
plt.ylabel('Departure delay (min)')
plt.title('Difficulty score vs Delay')
plt.tight_layout()
plt.savefig(plot_dir / 'difficulty_vs_delay.png')
plt.close()

# 3. Top destinations by avg difficulty
if 'scheduled_arrival_station_code' in score_df.columns:
    top_dests = score_df.groupby('scheduled_arrival_station_code')['difficulty_score'].mean().sort_values(ascending=False).head(10)
    plt.figure(figsize=(10,5))
    top_dests.plot(kind='bar')
    plt.ylabel('Avg difficulty score')
    plt.title('Top 10 destinations by avg difficulty')
    plt.tight_layout()
    plt.savefig(plot_dir / 'top_destinations.png')
    plt.close()
else:
    top_dests = pd.Series()

# 4. Boxplot of departure_delay by class
classes = score_df.groupby('difficulty_class')['departure_delay_minutes'].apply(list)
plt.figure(figsize=(8,4))
plt.boxplot([classes.get('Difficult',[]), classes.get('Medium',[]), classes.get('Easy',[])], labels=['Difficult','Medium','Easy'])
plt.ylabel('Departure delay (min)')
plt.title('Delay distribution by difficulty class')
plt.tight_layout()
plt.savefig(plot_dir / 'delay_by_class.png')
plt.close()

print('Saved plots to', plot_dir)

# === Regression quick test (SSR effect controlling for load) ===
try:
    import statsmodels.api as sm
    df_reg = score_df[['departure_delay_minutes','ssr_count','total_pax_flight']].dropna()
    X = df_reg[['ssr_count','total_pax_flight']]
    X = sm.add_constant(X)
    y = df_reg['departure_delay_minutes']
    model = sm.OLS(y, X).fit()
    reg_summary = model.summary().as_text()
    print('\nOLS regression (delay ~ ssr_count + total_pax_flight)')
    print(reg_summary)
except Exception as e:
    print('Statsmodels regression failed:', e)

# === Create PPTX with slides ===
prs = Presentation()
# basic slide width/height defaults are fine

# Helper to add title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'Flight Difficulty Score — ORD (sample)'
subtitle.text = 'Automated pipeline: EDA, Feature Engineering, Daily Difficulty Score, and Recommendations'

# Slide 2: Executive summary
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'Executive summary'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Goal: Create a daily Flight Difficulty Score to identify operations that need proactive resourcing.'
p = body.add_paragraph(); p.text = 'Output: CSV with difficulty_score (0-100), daily_rank, and class (Difficult/Medium/Easy).'
p.level = 1
p = body.add_paragraph(); p.text = 'Key drivers: low ground time, SSRs, transfer/hot-transfer bags, passenger load.'
p.level = 1

# Slide 3: Data & method
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Data & Method'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Data: Flight-level, PNR, PNR remarks, Bags, Airports (two weeks of ORD departures).'
p = body.add_paragraph(); p.text = 'Method: Feature engineering -> robust scaling -> weighted sum (weights from abs(correlation with delay) + floor) -> daily normalization (0-100) -> rank & classify.'

# Slide 4: EDA highlights + plot
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'EDA highlights'
left = Inches(0.5); top = Inches(1.2);
pic_path = plot_dir / 'difficulty_hist.png'
if pic_path.exists():
    slide.shapes.add_picture(str(pic_path), left, top, width=Inches(9))
# add small bullet textbox
tx = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
tf = tx.text_frame
tf.text = 'Avg delay: check CSV. Use distribution to choose thresholds for interventions.'

# Slide 5: Difficulty vs Delay
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Difficulty vs Delay'
pic_path = plot_dir / 'difficulty_vs_delay.png'
if pic_path.exists():
    slide.shapes.add_picture(str(pic_path), Inches(0.5), Inches(1.2), width=Inches(9))

# Slide 6: Top destinations
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Top destinations by avg difficulty'
pic_path = plot_dir / 'top_destinations.png'
if pic_path.exists():
    slide.shapes.add_picture(str(pic_path), Inches(0.5), Inches(1.2), width=Inches(9))
else:
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(2))
    tb.text_frame.text = 'No arrival station code available in data to compute top destinations.'

# Slide 7: Operational recommendations
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Operational recommendations'
body = slide.shapes.placeholders[1].text_frame
body.text = '1. Flag top 20% daily as Difficult and pre-assign resources.'
p = body.add_paragraph(); p.text = '2. Add dedicated CSR/agent for flights with high SSR counts.'; p.level = 1
p = body.add_paragraph(); p.text = '3. Prioritize baggage handling for flights with hot-transfer bags.'; p.level = 1
p = body.add_paragraph(); p.text = '4. Consider schedule buffers for flights with ground_time_deficit <= 5.'; p.level = 1

# Slide 8: Appendix / files
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Appendix & output files'
body = slide.shapes.placeholders[1].text_frame
body.text = f'CSV output: {out_csv}\nPPTX: Flight_Difficulty_Presentation.pptx\nPlots: {plot_dir}'

# Save PPTX
pptx_path = WORKDIR / 'Flight_Difficulty_Presentation.pptx'
prs.save(pptx_path)
print('Saved PPTX to', pptx_path)

# Final prints and optional download (in Colab you can download from the Files sidebar)
print('\nDone. Outputs:')
print(' - CSV:', out_csv)
print(' - PPTX:', pptx_path)
print(' - Plots dir:', plot_dir)

if COLAB:
    try:
        from google.colab import files
        print('\nYou can download files from the left Files panel, or use:')
        print("from google.colab import files; files.download('test_yourname.csv')")
    except Exception:
        pass

# ======= Script end =======
